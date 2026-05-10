from pptx import Presentation
import os

def add_slide(prs, title, bullets, notes, is_title=False):
    layout = prs.slide_layouts[0] if is_title else prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    
    if is_title:
        slide.shapes.title.text = bullets[0]
        slide.placeholders[1].text = bullets[1]
    else:
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.space_after = 6
            
    slide.notes_slide.notes_text_frame.text = notes

def main():
    prs = Presentation()
    
    # STRICT 3-TUPLE STRUCTURE: (title, bullets_list, trainer_notes)
    slides = [
        ("", ["Functions & Advanced Python",
              "Writing Reusable, Modular & Pythonic Automation Code\nPresenter: [Your Name]"],
         "🎯 OPENING: Ask 'How many of you have copied-pasted the same 10-line script across projects?' Set the tone: Functions aren't just syntax—they're the backbone of maintainable automation frameworks. Emphasize that today's session directly translates to pytest fixtures, API clients, and CI utilities."),
        
        ("Session Agenda & Learning Goals",
         ["• Function anatomy: parameters, returns, type hints",
          "• Flexible arguments: *args & **kwargs",
          "• Lambda expressions & when to use them",
          "• Functional programming: map, filter, reduce",
          "• Code modularization & framework-ready patterns",
          "• Hands-on: Build a configurable test data processor",
          "🎯 Outcome: Write cleaner, reusable, framework-compatible code"],
         "⏱️ PACING: 60 mins total. Keep theory tight; spend 40% on live coding & 20% on hands-on. Highlight that each concept maps directly to test framework design. Ask participants to open their IDEs before Slide 4."),
        
        ("Why Functions? The Automation Multiplier",
         ["• DRY Principle: Don't Repeat Yourself",
          "• Encapsulate logic → isolate test cases from setup",
          "• Easier debugging, testing, and CI/CD integration",
          "Example:",
          "  def get_api_token(env):",
          "      return requests.post(f'{env}/auth').json()['token']",
          "💡 Used in: Login flows, config loaders, data generators"],
         "🎤 ENGAGEMENT: Show a 30-line script with duplicated login logic. Refactor it into a function live. Ask: 'What breaks if the auth endpoint changes in 5 files vs 1 function?' Connect to framework maintainability."),
        
        ("Parameters & Return Types",
         ["• Positional, keyword, and mixed parameters",
          "• Type hints improve readability & IDE autocomplete",
          "• Multiple returns via tuples",
          "Example:",
          "  def parse_log(line: str) -> tuple[str, int]:",
          "      status = line.split()[0]",
          "      code = int(line.split()[1])",
          "      return status, code"],
         "⚠️ PITFALL: Learners ignore type hints. Stress that hints don't enforce types but enable linters (mypy, ruff) and prevent silent CI failures. Show IDE popup with `-> tuple[str, int]`. Relate to test data validation contracts."),
        
        ("Default & Keyword Arguments",
         ["• Defaults make functions flexible without breaking calls",
          "• Keyword args clarify intent in complex calls",
          "Example:",
          "  def setup_browser(headless=True, timeout=30, retries=3):",
          "      print(f'Config: headless={headless}, timeout={timeout}')",
          "💡 automation: Environment-specific overrides without rewriting"],
         "🎤 ENGAGEMENT: Run `setup_browser()` vs `setup_browser(timeout=10)`. Ask: 'Why are defaults dangerous if they're mutable?' (Preview: lists/dicts). Show how keyword args make test steps self-documenting in CI logs."),
        
        ("*args: Variable Positional Arguments",
         ["• Collects extra positional args into a tuple",
          "• Ideal for unknown-length input sequences",
          "Example:",
          "  def log_failures(*test_ids):",
          "      for tid in test_ids:",
          "          print(f'❌ Test {tid} failed')",
          "  log_failures(101, 102, 105)"],
         "💡 AUTOMATION CONTEXT: Bulk test tagging, retry queues, error aggregation. Show how `args` prevents hardcoding test counts. Warn against mixing `*args` with positional-only params without care. Live demo: pass 2 vs 10 IDs."),
        
        ("**kwargs: Variable Keyword Arguments",
         ["• Collects extra keyword args into a dict",
          "• Perfect for dynamic configs & API payloads",
          "Example:",
          "  def send_report(**metadata):",
          "      subject = metadata.get('subject', 'Test Run')",
          "      print(f'Sending: {subject} | {len(metadata)} fields')"],
         "🎯 USE CASE: CI pipeline metadata, environment variables, API headers. Emphasize `.get()` for safe access. Show how frameworks pass `**kwargs` through layers (test → driver → API) without rewriting signatures."),
        
        ("Combining *args & **kwargs",
         ["• Signature order: def func(pos, *args, **kwargs):",
          "• Powerful for CLI tools & flexible test runners",
          "Example:",
          "  def run_tests(*tags, **filters):",
          "      print(f'Tags: {tags} | Filters: {filters}')",
          "  run_tests('smoke', 'regression', env='prod', priority='high')"],
         "⚠️ ORDER MATTERS: Positional → *args → keyword → **kwargs. Break it intentionally in demo to show `SyntaxError`. Relate to pytest's `-k` and `--tags` patterns. Show how this scales to custom test runners."),
        
        ("Lambda Expressions",
         ["• Anonymous, single-expression functions",
          "• Best for short-lived operations (sorting, filtering)",
          "Example:",
          "  logs = [{'time': 12.3, 'msg': 'ok'}, {'time': 5.1, 'msg': 'fail'}]",
          "  sorted_logs = sorted(logs, key=lambda x: x['time'])"],
         "🎤 ENGAGEMENT: Ask 'When should you NOT use lambda?' (Complex logic, debugging, PEP-8 compliance). Show how `key=` replaces verbose `def` in test data sorting. Connect to pandas `apply()` and custom pytest markers."),
        
        ("Functional Programming: map() & filter()",
         ["• Transform & clean data without explicit loops",
          "• map(func, iterable) → applies func to each item",
          "• filter(func, iterable) → keeps items where func is True",
          "Example:",
          "  raw = ['  PASS ', 'FAIL', ' SKIP ']",
          "  clean = list(map(str.strip, raw))",
          "  fails = list(filter(lambda s: s == 'FAIL', clean))"],
         "💡 AUTOMATION CONTEXT: Log parsing, status normalization, payload sanitization. Emphasize that `map`/`filter` return iterators (memory efficient for large test suites). Live demo: chain them vs list comprehensions."),
        
        ("Functional Programming: reduce() & Comprehensions",
         ["• functools.reduce() accumulates values across sequence",
          "• Example: total_time = reduce(add, execution_times)",
          "• 🐍 Pythonic Alternative: List/Dict comprehensions",
          "  [x.upper() for x in logs if x['status'] == 'FAIL']",
          "• Prefer comprehensions for readability in test frameworks"],
         "⚠️ REALITY CHECK: `reduce()` is rarely used in modern Python. Show equivalent `sum()` or comprehension. Stress that test frameworks prioritize readability over FP purity. Recommend comprehensions for 90% of data transformations."),
        
        ("Code Modularization",
         ["• Split monolithic scripts into focused modules",
          "• Benefits: parallel CI runs, easier code review, reusability",
          "• Structure:",
          "  📁 automation/",
          "    ├── config.py",
          "    ├── utils.py",
          "    └── test_runner.py",
          "• Import: from utils import validate_payload"],
         "🎯 FRAMEWORK DESIGN: Show how pytest/Playwright organize code. Explain that modularization prevents 'spaghetti test suites'. Demo: move a function to `utils.py`, import it, run. Highlight CI pipeline caching benefits."),
        
        ("The __name__ == '__main__' Guard",
         ["• Prevents code execution when file is imported",
          "• Critical for test frameworks & reusable modules",
          "Example:",
          "  def run(): print('Starting automation...')",
          "  if __name__ == '__main__':",
          "      run()",
          "🚫 Without it: pytest imports → runs main logic unintentionally"],
         "⚠️ CRITICAL FOR AUTOMATION: Demonstrate what happens when you import a script without the guard. Show pytest output vs direct execution. Stress that all framework entry points MUST use this pattern to avoid side effects in CI."),
        
        ("Hands-On Challenge: Flexible Report Generator",
         ["Task: Build `generate_report(*tests, **config)`",
          "1. Accept variable test names & optional config (format, dest)",
          "2. Filter only 'passed' tests using filter() or comprehension",
          "3. Use lambda to sort by test name",
          "4. Print formatted summary based on config['format']",
          "🎯 Goal: Combine args, kwargs, FP, and modular thinking"],
         "⏱️ TIMING: 15 mins coding, 10 mins review. Provide starter scaffold. Walk through solutions emphasizing defensive `.get()`, type hints, and guard clause for missing config. Relate directly to Allure/pytest-html reporting plugins."),
        
        ("Best Practices & Framework Patterns",
         ["✅ Use type hints + docstrings for team collaboration",
         "✅ Keep functions pure (no side effects, return new data)",
         "✅ Prefer comprehensions over map/filter for readability",
         "✅ Use **kwargs for config layers, *args for queues/tags",
         "✅ Structure: utils/ → helpers, core/ → logic, tests/ → cases",
         "🔜 Next: pytest fixtures, parameterized tests, CI pipelines"],
         "🎯 DELIVERY: Summarize as 'The 5 Rules of Framework-Ready Functions'. Show how these practices reduce flaky tests and speed up PR reviews. Recommend `pylint`/`ruff` + `mypy` for automation repos. Open floor for real pain points."),
        
        ("Q&A, Resources & Next Steps",
         ["📖 Docs: docs.python.org/3/tutorial/controlflow.html",
          "🛠️ Practice: Exercism Python Track, Real Python Functions Guide",
          "📦 Starter Stack: pytest, pydantic (validation), rich (logging)",
          "🔜 Next Session: File I/O, pytest architecture, API clients",
          "💬 Questions? GitHub repo with slides + solutions shared"],
         "🎤 CLOSE: Thank participants. Share QR/code for repo. Ask 'Which concept will you refactor first in your current scripts?' Collect feedback. Preview next session's CI/CD integration focus. End on actionable note.")
    ]

    for i, (title, bullets, notes) in enumerate(slides):
        add_slide(prs, title, bullets, notes, is_title=(i == 0))

    filename = "Python_Functions_Advanced.pptx"
    prs.save(filename)
    print(f"✅ Successfully generated: {os.path.abspath(filename)}")

if __name__ == "__main__":
    main()