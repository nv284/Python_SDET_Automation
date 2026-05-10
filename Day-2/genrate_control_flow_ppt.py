from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def add_slide(prs, title, bullets, notes, is_title=False):
    layout = prs.slide_layouts[0] if is_title else prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    
    if is_title:
        slide.shapes.title.text = bullets[0]
        slide.placeholders[1].text = bullets[1]
    else:
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.space_after = Pt(6)
            
    slide.notes_slide.notes_text_frame.text = notes

def main():
    prs = Presentation()
    
    # CONSISTENT STRUCTURE: (title, bullets_list, trainer_notes)
    slides = [
        # Slide 1: Title Slide
        ("", ["Control Flow & Data Structures in Python", 
              "Mastering Logic, Collections & Iteration for Automation\nPresenter: [Your Name]"],
         "Welcome participants. Set expectations: This session transitions from basics to real automation patterns. Emphasize that control flow + collections form 80% of test framework logic."),
        
        # Slide 2
        ("Session Agenda",
         ["• Conditional logic (if, elif, else)",
          "• Nested conditions & real-time validation flows",
          "• Loops: for, while, break, continue",
          "• Collections: Lists (CRUD), Tuples, Sets, Dictionaries",
          "• Iteration patterns: enumerate(), zip()",
          "• Hands-on: Integrated data processing scenario",
          "• Best practices for automation frameworks"],
         "Walk through the flow: Logic → Collections → Iteration → Integration. Highlight that each concept will be tied to test data handling, API response parsing, and CI/CD logging."),
        
        # Slide 3
        ("Conditional Logic (if, elif, else)",
         ["• Branch execution based on boolean evaluation",
          "• Truthy/Falsy: non-empty strings, non-zero numbers → True; empty/0/None → False",
          "• Ternary operator: compact single-line assignments",
          "  status = 200",
          "  result = 'Pass' if status == 200 else 'Fail'",
          "• Automation tip: Prefer explicit comparisons over implicit truthiness"],
         "Explain Python's truthiness model. Show how `if response:` works but warn that `0` or `[]` evaluate False unintentionally. Emphasize readability: `if status_code == 200` > `if status_code:`."),
        
        # Slide 4
        ("Nested Conditions & Validation Flows",
         ["• Real-time validation often requires multiple checks",
          "• Deep nesting → 'Pyramid of Doom' → hard to maintain",
          "• Best practice: Guard clauses / early returns",
          "  def validate(payload):",
          "      if not payload: return False",
          "      if 'id' not in payload: return False",
          "      return process(payload)"],
         "Demonstrate how nested `if`s become unreadable in test frameworks. Show the guard clause pattern. Explain that early exits reduce cognitive load and make test logs cleaner."),
        
        # Slide 5
        ("Loops: for & while",
         ["• for loop: iterate over known sequences (lists, ranges, dicts)",
          "• while loop: repeat until a condition becomes False",
          "• for is preferred in automation (deterministic data)",
          "  for attempt in range(1, 4):",
          "      print(f'Retry {attempt}')",
          "• while = polling, retry logic, dynamic wait conditions"],
         "Clarify when to use each. `for` for test suites, CSV rows, API batches. `while` for retry mechanisms or waiting for UI elements. Warn about infinite `while` loops without timeouts in CI."),
        
        # Slide 6
        ("Loop Control: break & continue",
         ["• break: exit loop immediately",
          "• continue: skip current iteration, jump to next",
          "• Prevents unnecessary processing & speeds up test runs",
          "  for test in suite:",
          "      if test.skipped: continue",
          "      if test.critical and test.failed: break"],
         "Show practical use: `continue` for skipping flaky tests. `break` for failing fast on critical blockers. Emphasize logging why a loop exits early to avoid hidden test failures."),
        
        # Slide 7
        ("Lists: CRUD Operations",
         ["• Ordered, mutable, allows duplicates",
          "• Create: lst = [] or [1, 2, 3]",
          "• Read: lst[0], lst[-1], lst[1:3]",
          "• Update: lst.append(), lst.insert(), lst[1] = 'new'",
          "• Delete: lst.pop(), lst.remove(), del lst[0]"],
         "Lists are the backbone of test data management. Demonstrate slicing (`lst[-2:]`) for batching. Warn about `remove()` throwing ValueError if item missing. Connect to test case queues."),
        
        # Slide 8
        ("Tuples: Immutability & Use Cases",
         ["• Ordered, IMMUTABLE, slightly faster than lists",
          "• Syntax: tup = ('dev', 8080) or tup = 1, 2, 3",
          "• Unpacking: host, port = tup",
          "• Use cases: config constants, DB coordinates, fixed headers"],
         "Explain why immutability matters: prevents accidental test state mutations in parallel runs. Show unpacking for clean config assignment. Note that tuples can hold mutable objects inside."),
        
        # Slide 9
        ("Sets: Unique Data Handling",
         ["• Unordered, stores ONLY unique elements",
          "• Fast O(1) lookups, ideal for deduplication",
          "• Math ops: union (|), intersection (&), difference (-)",
          "  executed = {101, 102, 101}  # → {101, 102}",
          "  missing = all_ids - executed"],
         "Sets shine in automation: tracking executed test IDs, finding missing test cases, or comparing expected vs actual results. Warn that sets are unordered; don't use them when sequence matters."),
        
        # Slide 10
        ("Dictionaries: Key-Value Mapping",
         ["• Hash map: fast lookups, JSON-like structure",
          "• Keys must be immutable (str, int, tuple)",
          "• Ideal for API responses, test configs, environment vars",
          "  resp = {'status': 200, 'data': {'id': 5}}",
          "  print(resp['status'])  # 200"],
         "Dictionaries are the default structure for modern APIs and test fixtures. Emphasize that keys are unique. Show how nested dicts mirror JSON payloads. Warn against assuming keys exist."),
        
        # Slide 11
        ("Dictionaries: Safe Access & Merging",
         ["• .get(key, default) → avoids KeyError crashes",
          "• .update() or {**d1, **d2} → merge configs safely",
          "• .keys(), .values(), .items() → iteration helpers",
          "  token = resp.get('token', 'expired')",
          "  final_cfg = {**defaults, **user_overrides}"],
         "Stress defensive programming: `.get()` prevents test suite crashes on malformed API responses. Merging is essential for environment-specific config overrides. Show how `.items()` pairs with loops."),
        
        # Slide 12
        ("Iteration Pattern: enumerate()",
         ["• Track index + value simultaneously",
          "• Replaces manual counters & range(len())",
          "• Cleaner logging & test numbering",
          "  for i, test in enumerate(cases, start=1):",
          "      print(f'{i}. Running {test}')"],
         "Show why `for i in range(len(list))` is anti-pattern. `enumerate()` is Pythonic, readable, and less error-prone. Relate to generating sequential test logs, numbering JUnit reports, or tracking retries."),
        
        # Slide 13
        ("Iteration Pattern: zip()",
         ["• Parallel iteration over multiple sequences",
          "• Stops at the shortest sequence length",
          "• Perfect for mapping headers ↔ rows, inputs ↔ expected",
          "  for col, val in zip(headers, row_data):",
          "      print(f'{col}: {val}')"],
         "Demonstrate CSV/API response parsing. Show how `zip` eliminates index math. Warn about silent truncation if sequences differ in length. Suggest `itertools.zip_longest` when alignment matters."),
        
        # Slide 14
        ("Real-World Automation Pattern",
         ["• Combine control flow + collections for data pipelines",
          "• Filter test cases → validate → execute → report",
          "  valid_tests = [t for t in suite if t['env'] == current_env]",
          "  for idx, test in enumerate(valid_tests, 1):",
          "      run(test) if test['enabled'] else skip(test)"],
         "Walk through how frameworks actually process data. List comprehensions + filtering + enumerated execution mirror pytest collection phases. Emphasize that readable pipelines reduce maintenance overhead."),
        
        # Slide 15
        ("Hands-On Challenge",
         ["Task: Process mock test results & generate summary",
          "1. Parse list of dicts: {'id': 1, 'status': 'pass', 'time': 1.2}",
          "2. Filter only 'fail' cases using list comprehension",
          "3. Use enumerate to number failures",
          "4. Calculate avg time using sum() & len()",
          "Goal: Apply all concepts in one cohesive script"],
         "Guide them step-by-step. Encourage pairing. Suggest using `.get()` for safe status access. After 10 mins, review a clean solution. Highlight how this mirrors real CI log parsing."),
        
        # Slide 16
        ("Summary & Automation Best Practices",
         ["✅ Guard clauses > deep nesting",
          "✅ for loops + enumerate/zip = Pythonic iteration",
          "✅ Choose collection by need: List(order), Set(unique), Dict(map), Tuple(constant)",
          "✅ .get() & validation prevent flaky automation",
          "🔜 Next: Functions, pytest fixtures, file I/O, API clients",
          "💬 Q&A & GitHub repo for all examples"],
         "Reinforce the 'right tool for the job' mindset. Recommend Black/Ruff for auto-formatting. Preview how these concepts scale to parameterized tests. Open floor for real-world scenarios.")
    ]

    for i, (title, bullets, notes) in enumerate(slides):
        add_slide(prs, title, bullets, notes, is_title=(i == 0))

    filename = "Python_ControlFlow_DataStructures.pptx"
    prs.save(filename)
    print(f"✅ Successfully generated: {os.path.abspath(filename)}")

if __name__ == "__main__":
    main()