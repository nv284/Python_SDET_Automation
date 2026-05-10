from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def add_slide(prs, title, bullets, notes, is_title_slide=False):
    if is_title_slide:
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = bullets[0]
    else:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.space_after = Pt(6)
            p.level = 0
    
    # Add trainer notes
    slide.notes_slide.notes_text_frame.text = notes

def main():
    prs = Presentation()
    
    # Slide Data: (Title, Bullets, Trainer Notes)
    slides = [
        ("title", [
            "Python Fundamentals for Automation",
            "Concepts, Examples & Hands-on Labs | Presenter: [Your Name]"
        ], "Welcome participants. Set expectations: This session bridges Python basics with real-world automation use cases. Emphasize that we'll focus on writing clean, maintainable scripts, not just syntax."),
        
        ("content", "Session Agenda", [
            "• Python overview for automation ecosystem",
            "• Installation & environment setup (venv, pip)",
            "• Variables & memory concepts",
            "• Core data types: int, float, str, bool",
            "• Type casting & validation techniques",
            "• Operators: arithmetic, relational, logical",
            "• Input/Output handling & formatting",
            "• Hands-on: Scripts, validation, string scenarios"
        ], "Walk through the agenda. Highlight the flow: theory → live demo → hands-on labs → Q&A. Mention that every concept will be tied back to test/data/CI automation. Ask participants about their current scripting experience."),
        
        ("content", "Python Overview for Automation Ecosystem", [
            "• Why Python dominates automation: readable syntax, massive ecosystem, cross-platform",
            "• Key libraries: requests (API), Selenium/Playwright (UI), pytest (testing), pandas (data)",
            "• Seamless CI/CD integration (GitHub Actions, Jenkins, GitLab CI)",
            "Example: import subprocess; subprocess.run([\"git\", \"status\"])",
            "• Scales from 5-line scripts to enterprise test frameworks"
        ], "Contrast with Java/C# verbosity. Mention Python's role in DevOps/test automation. Ask: 'What automation tasks do you currently do manually?' Highlight that Python's 'batteries-included' philosophy speeds up prototyping and productionization."),
        
        ("content", "Installation & Environment Setup (venv, pip)", [
            "• Global installs cause dependency conflicts → use virtual environments",
            "• venv: isolated Python environment per project",
            "• pip: package manager for installing libraries",
            "Commands:",
            "  python -m venv .venv",
            "  source .venv/bin/activate  # Mac/Linux",
            "  .venv\\Scripts\\activate    # Windows",
            "  pip install requests pytest",
            "  pip freeze > requirements.txt"
        ], "Emphasize that professional projects ALWAYS use venvs. Show Windows vs Mac/Linux activation differences. Warn against pip install without venv in shared systems. Explain how requirements.txt ensures team reproducibility and CI pipeline stability."),
        
        ("content", "Variables & Memory Concepts", [
            "• Variables are REFERENCES to objects, not storage boxes",
            "• Dynamic typing: type is determined at runtime",
            "• Immutable (int, str, tuple) vs Mutable (list, dict) affects memory behavior",
            "Example: a = 10; b = a; a = 20; print(b) → 10",
            "• id() shows memory address; useful for debugging test data leaks"
        ], "Use the 'name tag' analogy. Demonstrate id() live. Explain why this matters in automation: passing test data by reference vs value, avoiding unintended state mutations in test suites. Show how lst.append() changes shared references."),
        
        ("content", "Core Data Types", [
            "• int: arbitrary precision whole numbers (42, -5)",
            "• float: IEEE 754 decimals (3.14, 0.0) → beware precision!",
            "• str: Unicode text, IMMUTABLE ('Hello', 'https://api.test')",
            "• bool: True/False + truthy/falsy evaluation (empty str/list → False)",
            "• type() vs isinstance() for robust type checking"
        ], "Highlight float precision issues (0.1 + 0.2 != 0.3). Explain why strings are immutable (safe for concurrent test runs). Relate to test data types (IDs, URLs, flags). Show truthy/falsy behavior in conditional assertions."),
        
        ("content", "Type Casting & Validation", [
            "• Explicit casting: int(), float(), str(), bool()",
            "• Validation: .isdigit(), .isalpha(), try/except ValueError",
            "• input() ALWAYS returns a string → casting without validation crashes scripts",
            "Example: val = int(user_input) → ValueError if non-numeric",
            "• Best practice: validate → cast → process"
        ], "Stress that input() always returns str. Show real automation failure: reading CSV/API response as string, trying math without cast. Introduce try/except as industry standard. Demonstrate how validation prevents CI pipeline failures."),
        
        ("content", "Operators (Arithmetic, Relational, Logical)", [
            "• Arithmetic: + - * / // % **",
            "• Relational: == != > < >= <= (returns bool)",
            "• Logical: and or not (short-circuit evaluation)",
            "• Operator precedence: parentheses recommended for readability",
            "• Automation use: assertions, conditional test flows, status checks"
        ], "Connect to test assertions. Explain short-circuiting (False and expensive_call() skips call). Warn about == vs is (value vs identity). Relate to conditional test flows like retry logic or environment-specific checks."),
        
        ("content", "Input/Output Handling & Formatting", [
            "• input(): reads CLI input as string",
            "• print(): modern formatting via f-strings (Python 3.6+)",
            "• f-string specifiers: :.2f (decimals), >10 (padding), etc.",
            "• Production automation: prefer logging module over print()",
            "Example: print(f'✅ Test passed in {t:.3f}s | Status: {s}')"
        ], "Show why f-strings win (speed, readability, inline expressions). Mention logging module for production automation (levels, file output, structured logs). Relate to CI/CD log parsing and debugging failed runs."),
        
        ("content", "Hands-on 1: Build Simple Scripts", [
            "Task: Rectangle Area & Perimeter Calculator",
            "Steps: take input → cast to float → calculate → format output",
            "Code: length = float(input('Length: ')); width = float(input('Width: '));",
            "      area = length * width; print(f'Area: {area:.2f}')",
            "Goal: Combine I/O, casting, arithmetic, formatting in one flow"
        ], "Guide them step-by-step. Ask them to run it. Then break it intentionally (enter text) to show need for validation. Discuss how this mirrors reading test config files or CSV data."),
        
        ("content", "Hands-on 2: Input Validation Program", [
            "Task: Robust Age/Number Validator (18-100 range)",
            "Pattern: while True → try/except → range check → break on success",
            "Code: while True: try: age = int(input('Age: ')); if 18<=age<=100: break",
            "      except ValueError: print('Numbers only!')",
            "Goal: Defensive programming for automation test data"
        ], "Explain defensive programming. Show how this pattern mirrors test data validation before API/UI interactions. Discuss ValueError handling, loop control, and user feedback. Relate to CI data validation gates."),
        
        ("content", "Hands-on 3: String-Based Test Scenarios", [
            "Task: Validate Username Rules (3-15 chars, alphanumeric, starts with letter)",
            "Methods: len(), .isalpha(), .isalnum(), .strip(), .lower()",
            "Code: def validate(u): return 3<=len(u)<=15 and u[0].isalpha() and u.isalnum()",
            "• String manipulation = 70% of UI/API test data prep",
            "• Normalize inputs before assertions for reliable tests"
        ], "Show how string manipulation is critical for test data prep. Connect to regex mention (but keep it simple). Emphasize idempotency in test data. Discuss case-insensitive comparisons and whitespace trimming in automation."),
        
        ("content", "Summary & Automation Best Practices", [
            "✅ Python is automation-ready out of the box",
            "✅ Environments keep dependencies isolated & reproducible",
            "✅ Variables reference memory; immutability prevents side effects",
            "✅ Validate → Cast → Process prevents runtime crashes",
            "✅ Use f-strings & logging for clean, parseable automation output"
        ], "Reinforce automation mindset: readable, reusable, validated. Recommend practice resources. Explain how these fundamentals scale to pytest fixtures, API clients, and CI scripts. Open floor for real-world scenarios."),
        
        ("content", "Q&A & Next Steps", [
            "📖 Official Docs: docs.python.org/3/",
            "🛠️ Practice: Exercism, Codewars (Python track), pytest docs",
            "📦 Starter Stack: VS Code + Python Ext + venv + Black/Ruff",
            "🔜 Next Session: Functions, File I/O, pytest, Selenium/Requests intro",
            "💬 Questions? Feedback? GitHub repo for all examples provided"
        ], "Collect feedback. Share GitHub repo link with all slide examples + solutions. Encourage pairing up for hands-on practice. Preview next session structure. Thank participants.")
    ]

    # Generate Slides
    for i, item in enumerate(slides):
        if i == 0:
            add_slide(prs, item[1][0], item[1][1:], item[2], is_title_slide=True)
        else:
            add_slide(prs, item[1], item[2], item[3])

    # Save
    filename = "Python_Automation_Fundamentals.pptx"
    prs.save(filename)
    print(f"✅ Successfully generated: {os.path.abspath(filename)}")

if __name__ == "__main__":
    main()