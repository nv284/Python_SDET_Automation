from pptx import Presentation
import os

def add_slide(prs, title, bullets, notes, is_title=False):
    layout = prs.slide_layouts[0] if is_title else prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    
    if is_title:
        slide.shapes.title.text = bullets[0]
        slide.placeholders[1].text = bullets[1]
    else:
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.space_after = 6
            
    slide.notes_slide.notes_text_frame.text = notes

def main():
    prs = Presentation()
    
    # STRICT 3-TUPLE STRUCTURE: (title, bullets_list, trainer_notes)
    slides = [
        ("", ["OOP for Automation Frameworks",
              "Building Scalable, Maintainable Test Architecture\nPresenter: [Your Name]"],
         "[GOAL] OPENING: Ask 'How many of your test scripts break when a single UI locator changes?' Set the tone: OOP isn't academic theory—it's the engineering foundation that separates fragile scripts from enterprise-grade automation frameworks. Emphasize Page Object Model (POM) and framework scalability from Day 1."),
        
        ("Session Agenda & Learning Goals",
         ["• Classes, objects & state management",
          "• Constructors & instance variables in test contexts",
          "• Inheritance: Reusing framework logic",
          "• Polymorphism: Environment-specific overrides",
          "• Encapsulation & Abstraction: Hiding complexity",
          "• Mapping OOP to Framework Design (POM intro)",
          "• Hands-on: Build a mini framework module",
          "[OUTCOME] Write framework-ready, team-scalable code"],
         "[TIME] PACING: 70 mins. 30% theory, 40% live refactoring, 30% hands-on. Highlight that every OOP concept maps directly to pytest/Selenium/Playwright architecture. Ask participants to open their IDEs before Slide 4."),
        
        ("Why OOP for Automation?",
         ["• Scripting vs Framework: Linear vs Modular",
          "• DRY Principle: Change locators once, not 50 times",
          "• Isolated state for parallel test execution",
          "• Cleaner CI/CD logs & easier code reviews",
          "Example: 30-line duplicate login → 1 reusable class method"],
         "[ASK] Show a messy script with repeated waits, locators, and assertions. Ask: 'What breaks first when QA changes the login button ID?' Connect to framework maintainability, flaky test reduction, and team onboarding speed."),
        
        ("Classes & Objects: Blueprints vs Instances",
         ["• Class = template (behavior + structure)",
          "• Object = live instance in memory",
          "• 'self' refers to the current instance",
          "Example:",
          "  class BrowserSession:",
          "      def navigate(self, url): pass",
          "  chrome = BrowserSession()  # Instance 1"],
         "[DEMO] Use the 'car factory' analogy but map to 'test runner'. Emphasize that `self` is just a reference to the specific object. Show `id(chrome)` vs `id(firefox)` to prove isolated memory spaces for parallel runs."),
        
        ("Constructors & Instance Variables",
         ["• __init__ initializes object state",
          "• Instance vars (self.driver) live per test",
          "• Prevents global state pollution",
          "Example:",
          "  class LoginPage:",
          "      def __init__(self, driver, base_url):",
          "          self.driver = driver"],
         "[WARN] Show the anti-pattern: global `driver = webdriver.Chrome()`. Explain why it causes flaky parallel tests. Instance variables = isolated test state. Demo `__init__` accepting config vs hardcoding."),
        
        ("Inheritance: Framework Code Reuse",
         ["• Child classes inherit parent methods",
          "• BasePage → LoginPage, DashboardPage",
          "• Centralize waits, screenshots, logging",
          "Example:",
          "  class BasePage:",
          "      def wait_for_element(self, loc): ...",
          "  class LoginPage(BasePage): ..."],
         "[TIP] Inheritance solves duplication but creates tight coupling. Show how `BasePage` holds 80% of framework utilities. Demo `super().__init__()` and warn against inheritance trees >2 levels deep."),
        
        ("Polymorphism: Same Interface, Different Behavior",
         ["• Override parent methods in child classes",
          "• Env-specific auth, region-specific flows",
          "• Tests call .login() → framework picks impl",
          "Example:",
          "  class SAMLPage(LoginPage):",
          "      def login(self, creds):  # override",
          "          self.click_idp_button()"],
         "[GOAL] Show how CI can swap implementations without touching test code. Polymorphism = 'write once, run anywhere' for automation. Demo pytest parametrization + polymorphic page objects."),
        
        ("Encapsulation: Hiding Implementation Details",
         ["• Keep locators private (_locators)",
          "• Expose clean API (search(), submit())",
          "• Tests don't care about XPath/CSS changes",
          "Example:",
          "  class SearchPage:",
          "      _SEARCH_INPUT = '#search'",
          "      def enter_query(self, q): ..."],
         "[ASK] 'Should tests know the button's XPath?' No. Encapsulation = UI changes don't break tests. Show `@property` for computed attributes. Relate to black-box testing principles in frameworks."),
        
        ("Abstraction: Enforcing Framework Contracts",
         ["• ABC (Abstract Base Class) defines rules",
          "• Forces child classes to implement methods",
          "• Prevents missing overrides in team code",
          "Example:",
          "  from abc import ABC, abstractmethod",
          "  class BaseDriver(ABC):",
          "      @abstractmethod",
          "      def click(self): ..."],
         "[WARN] Abstraction catches missing implementations at import time, not at runtime. Show how it standardizes team contributions. Critical for multi-dev frameworks and plugin architectures."),
        
        ("Mapping OOP to Framework Design (POM)",
         ["• Page Object Model = industry standard",
          "• UI elements = attributes, Actions = methods",
          "• Tests = orchestration, Pages = implementation",
          "Example:",
          "  page = LoginPage(driver)",
          "  page.login(user)",
          "  assert page.is_dashboard()"],
         "[DEMO] Draw the separation: Tests call page methods, pages interact with UI. Emphasize NO assertions in page classes (unless returning state). Show how this scales to 100+ tests with 10 page files."),
        
        ("Framework Architecture: Base Classes & Mixins",
         ["• Base classes = shared behavior",
          "• Mixins = cross-cutting features (logging, retry)",
          "• Composition > deep inheritance",
          "Example:",
          "  class RetryMixin: def retry(self, func): ...",
          "  class LoginPage(BasePage, RetryMixin): ..."],
         "[TIP] Modern frameworks favor composition. Mixins add capabilities without bloating inheritance. Show how pytest fixtures + mixins create flexible, testable architectures. Warn against God classes."),
        
        ("Factory Pattern & Config Management",
         ["• Decouple object creation from usage",
          "• BrowserFactory.get(env, headless)",
          "• CI/CD driven setup, zero code changes",
          "Example:",
          "  driver = BrowserFactory.create('chrome')",
          "  page = LoginPage(driver)"],
         "[GOAL] Factories enable cloud execution (Sauce, BrowserStack) with 1-line config changes. Show how `**config` dicts flow from CI → factory → pages. Relate to environment matrix testing."),
        
        ("Hands-On Challenge: Mini POM Module",
         ["Task: Build BasePage + SearchPage",
          "1. Define BasePage with _driver, wait_for()",
          "2. Create SearchPage inheriting BasePage",
          "3. Encapsulate locators, expose search(query)",
          "4. Instantiate & run mock search",
          "[OUTCOME] Framework-ready, team-scalable module"],
         "[TIME] 15 mins coding, 10 mins review. Provide scaffold. Walk through emphasizing encapsulation, inheritance, and clean API. Relate directly to Selenium/Playwright page structure. Peer review encouraged."),
        
        ("Best Practices & Anti-Patterns",
         ["✅ Favor composition over deep inheritance",
          "✅ Keep pages thin (no assertions in pages)",
          "✅ Use properties for computed state",
          "❌ Avoid God objects, global drivers, mixed data/logic",
          "❌ Don't expose locators to test files"],
         "[WARN] Show common framework failures: pages doing DB queries, tests knowing CSS, massive Base classes. Recommend `pytest` + `page-objects` separation. Stress code review checklists."),
        
        ("Real-World Framework Mapping",
         ["• pytest fixtures + POM = scalable test runs",
          "• Playwright: Page class as first-class citizen",
          "• Selenium: PageFactory (deprecated) → manual POM",
          "• CI: Parallel execution via isolated instances"],
         "[TIP] Modern tools are OOP-native. Playwright's `Page`, Selenium's `WebDriver`, pytest's `Fixture` all follow OOP principles. Show how cloud grids scale instances, not scripts."),
        
        ("Summary & Framework Mindset",
         ["✅ OOP = maintainable, scalable, team-friendly",
          "✅ POM decouples tests from UI volatility",
          "✅ Encapsulation + Abstraction = stable contracts",
          "✅ Inheritance + Composition = flexible architecture",
          "[NEXT] CI/CD pipelines, parallel execution, reporting"],
         "[GOAL] Reinforce that OOP isn't about syntax—it's about engineering resilience. Frameworks that scale treat UI as an implementation detail, not a test dependency. Open floor for real architecture pain points."),
        
        ("Q&A, Resources & Next Steps",
         ["📖 Docs: docs.python.org/3/tutorial/classes.html",
          "🛠️ Practice: Test Automation University, Selenium POM Guide",
          "📦 Starter Stack: pytest, playwright, pydantic, allure",
          "[NEXT] CI/CD integration, parallel test execution, dashboard reporting",
          "[ASK] Questions? Which OOP principle will you refactor first?"],
         "[CLOSE] Thank participants. Share repo link. Preview next session's CI/CD focus. Ask for one actionable takeaway. End on engineering mindset, not just automation execution.")
    ]

    for i, (title, bullets, notes) in enumerate(slides):
        add_slide(prs, title, bullets, notes, is_title=(i == 0))

    filename = "Python_OOP_for_Automation.pptx"
    prs.save(filename)
    print(f"✅ Successfully generated: {os.path.abspath(filename)}")

if __name__ == "__main__":
    main()